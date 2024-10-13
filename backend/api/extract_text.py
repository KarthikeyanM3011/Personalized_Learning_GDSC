from flask import Blueprint, request, jsonify
import mimetypes
import io
import PyPDF2
from docx import Document
from pptx import Presentation

upload_bp = Blueprint('upload', __name__)

mimetypes.add_type('application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.docx')
mimetypes.add_type('application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx')

@upload_bp.route('/upload_files', methods=['POST'])
def upload_files():
    try:
        if 'files' not in request.files:
            return jsonify({"message": "No files uploaded"}), 400

        files = request.files.getlist('files')
        response_data = []

        for file in files:
            file_name = file.filename
            mime_type, _ = mimetypes.guess_type(file_name)

            if mime_type == 'application/pdf':
                text = extract_pdf_content(file)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                text = extract_word_content(file)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                text = extract_ppt_content(file)
            else:
                return jsonify({"message": f"Unsupported file type: {file_name}"}), 400
            
            response_data.append({
                "filename": file_name,
                "content": text
            })

        return jsonify({
            "message": "Files processed successfully",
            "files": response_data
        }), 200

    except Exception as e:
        print(e)
        return jsonify({"message": str(e)}), 500


def extract_pdf_content(file):
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text() + '\n'
        return text.strip()

    except Exception as e:
        print(e)
        return f"Error extracting PDF content: {str(e)}"


def extract_word_content(file):
    try:
        doc = Document(file)
        text = ''
        for para in doc.paragraphs:
            text += para.text + '\n'
        return text.strip()

    except Exception as e:
        print(e)
        return f"Error extracting Word content: {str(e)}"


def extract_ppt_content(file):
    try:
        presentation = Presentation(file)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text.strip()

    except Exception as e:
        print(e)
        return f"Error extracting PPT content: {str(e)}"
